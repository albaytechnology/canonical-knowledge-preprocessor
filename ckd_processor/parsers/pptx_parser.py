"""
Presentation Parser (.pptx, .ppt) extracting slide titles, body shapes, and notes.
"""

import os
from typing import List

from ckd_processor.parsers.base import BaseParser, ParsedDocument, PageContent
from ckd_processor.utils.logger import setup_logger

logger = setup_logger("PPTXParser")

try:
    import pptx
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class PPTXParser(BaseParser):
    def supported_extensions(self) -> List[str]:
        return ["pptx", "ppt"]

    def parse(self, filepath: str) -> ParsedDocument:
        filename = os.path.basename(filepath)
        pages: List[PageContent] = []

        if not HAS_PPTX:
            pages.append(PageContent(page_number=1, text="[pptx library missing. Please install python-pptx]"))
            return ParsedDocument(
                filepath=filepath,
                filename=filename,
                extension=filepath.split(".")[-1].lower(),
                pages=pages,
                full_raw_text=pages[0].text
            )

        try:
            prs = pptx.Presentation(filepath)
            for idx, slide in enumerate(prs.slides, start=1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())

                slide_content = f"### Slide {idx}\n\n" + "\n\n".join(slide_texts)
                pages.append(PageContent(page_number=idx, text=slide_content))

        except Exception as e:
            logger.error(f"Error parsing PPTX file {filename}: {e}")
            pages.append(PageContent(page_number=1, text=f"[PPTX Parse Error: {e}]"))

        full_text = "\n\n".join([p.text for p in pages])
        return ParsedDocument(
            filepath=filepath,
            filename=filename,
            extension=filepath.split(".")[-1].lower(),
            pages=pages,
            full_raw_text=full_text,
            raw_metadata={"slide_count": len(pages)}
        )
